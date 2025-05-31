import os
import threading
import io

import streamlit as st

from fastapi import FastAPI, File, UploadFile
from fastapi.responses import JSONResponse
import uvicorn

import pandas as pd
import docx
import pptx
import pytesseract
from PIL import Image
import pdfplumber

# ---- Import your utils ----
from utils.summarizer import gpt_generate_query, gpt_summarize
from utils.browser import google_search_and_scrape
from utils.extractor import extract_main_text

# ---------------------------
# Shared function: extract text from any supported file
def extract_text_from_file(file_bytes, filetype, filename=None):
    try:
        if filetype == "application/pdf":
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                return "\n".join([page.extract_text() or "" for page in pdf.pages])
        elif filetype in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"]:
            doc = docx.Document(io.BytesIO(file_bytes))
            return "\n".join([p.text for p in doc.paragraphs])
        elif filetype in [
            "application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]:
            df = pd.read_excel(io.BytesIO(file_bytes))
            return df.to_string()
        elif filetype in [
            "application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text])
        elif filetype in ["image/jpeg", "image/png"]:
            image = Image.open(io.BytesIO(file_bytes))
            return pytesseract.image_to_string(image)
        else:
            return "Unsupported file type."
    except Exception as e:
        return f"Extraction failed: {str(e)}"

# ---------------------------
# FastAPI server for external (Custom GPT) integration
app_api = FastAPI()

@app_api.post("/research")
async def research_file(file: UploadFile = File(...)):
    file_bytes = await file.read()
    text = extract_text_from_file(file_bytes, file.content_type, file.filename)
    query = gpt_generate_query(text[:2000])
    html = google_search_and_scrape(query)
    web_article = extract_main_text(html)
    crosscheck_prompt = (
        f"Uploaded Document Content:\n{text[:1000]}\n\n---\n"
        f"Web Source Content:\n{web_article[:2000]}\n\n"
        "Compare and report agreements, disagreements, or discrepancies. Summarize the main points."
    )
    summary = gpt_summarize(crosscheck_prompt)
    return JSONResponse({"summary": summary, "query": query})

def run_fastapi():
    uvicorn.run(app_api, host="0.0.0.0", port=8000)

# ---------------------------
# Start FastAPI in background
threading.Thread(target=run_fastapi, daemon=True).start()

# ---------------------------
# Streamlit UI (for local research)
st.set_page_config(page_title="AI Multimodal Search & Research", layout="wide")
st.title("🔎 AI Multimodal Search & Document Researcher")

uploaded = st.file_uploader(
    "Upload a document (DOCX, XLSX, PPTX, PDF, JPG, PNG)",
    type=["pdf", "docx", "xlsx", "pptx", "jpg", "jpeg", "png"])

if uploaded:
    file_bytes = uploaded.read()
    text = extract_text_from_file(file_bytes, uploaded.type, uploaded.name)
    st.text_area("📄 Extracted Content", text[:4000], height=300)

    if st.button("Research & Cross-Check Online"):
        with st.spinner("Generating research query from document..."):
            query = gpt_generate_query(text[:2000])
            st.write("**Generated Search Query:**", query)
        with st.spinner("Searching Google and extracting web content..."):
            html = google_search_and_scrape(query)
            web_article = extract_main_text(html)
        with st.spinner("LLM: Comparing Document & Online Source..."):
            crosscheck_prompt = (
                f"Uploaded Document Content:\n{text[:1000]}\n\n---\n"
                f"Web Source Content:\n{web_article[:2000]}\n\n"
                "Compare and report agreements, disagreements, or discrepancies. Summarize the main points and note any important findings."
            )
            summary = gpt_summarize(crosscheck_prompt)
            st.success(summary)

st.markdown("""
---
**Supports**: PDF, DOCX, XLSX, PPTX, JPG, PNG  
- Uses OCR for images  
- LLM-powered summarization & fact cross-check  
- Requires [Tesseract OCR](https://github.com/tesseract-ocr/tesseract) for image support  
- FastAPI endpoint available at `/research` (for Custom GPT or external API use)
""")

from fastapi import FastAPI, File, UploadFile
from fastapi.responses import JSONResponse
from mangum import Mangum
import io, os
import pandas as pd
import docx
import pptx
import pytesseract
from PIL import Image
import pdfplumber

# If your utils are in the root-level 'utils' folder
import sys
sys.path.append(os.path.join(os.path.dirname(__file__), '..'))
from utils.summarizer import gpt_generate_query, gpt_summarize
from utils.browser import google_search_and_scrape
from utils.extractor import extract_main_text

app = FastAPI()

def extract_text_from_file(file_bytes, filetype, filename=None):
    try:
        if filetype == "application/pdf":
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                return "\n".join([page.extract_text() or "" for page in pdf.pages])
        elif filetype in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"]:
            doc = docx.Document(io.BytesIO(file_bytes))
            return "\n".join([p.text for p in doc.paragraphs])
        elif filetype in [
            "application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]:
            df = pd.read_excel(io.BytesIO(file_bytes))
            return df.to_string()
        elif filetype in [
            "application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text])
        elif filetype in ["image/jpeg", "image/png"]:
            image = Image.open(io.BytesIO(file_bytes))
            return pytesseract.image_to_string(image)
        else:
            return "Unsupported file type."
    except Exception as e:
        return f"Extraction failed: {str(e)}"

@app.post("/research")
async def research_file(file: UploadFile = File(...)):
    file_bytes = await file.read()
    text = extract_text_from_file(file_bytes, file.content_type, file.filename)
    query = gpt_generate_query(text[:2000])
    html = google_search_and_scrape(query)
    web_article = extract_main_text(html)
    crosscheck_prompt = (
        f"Uploaded Document Content:\n{text[:1000]}\n\n---\n"
        f"Web Source Content:\n{web_article[:2000]}\n\n"
        "Compare and report agreements, disagreements, or discrepancies. Summarize the main points."
    )
    summary = gpt_summarize(crosscheck_prompt)
    return JSONResponse({"summary": summary, "query": query})

handler = Mangum(app)  # <-- Required by Vercel to serve the FastAPI app!
